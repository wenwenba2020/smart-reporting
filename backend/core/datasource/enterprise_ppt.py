"""
EnterprisePPTAdapter — parses .pptx files into structured ContentDeck with
LLM-powered slide summarization and style-profile extraction.
"""

from __future__ import annotations

import json
import logging
import uuid
from pathlib import Path
from typing import Optional

from backend.core.datasource.base import DataSourceAdapter, register_adapter
from backend.core.models import (
    ContentDeck,
    SlideAsset,
    SourceDocument,
    StyleProfile,
)

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Helpers (module-private)
# ---------------------------------------------------------------------------


def _safe_get(obj: object, attr: str, default=None):
    """Safely get an attribute that may not exist on python-pptx objects."""
    try:
        return getattr(obj, attr, default)
    except Exception:
        return default


def _extract_text_from_shape(shape) -> str:
    """Extract all text from a shape (text frame, table, group)."""
    texts: list[str] = []

    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t:
                texts.append(t)

    if shape.has_table:
        table = shape.table
        for row in table.rows:
            row_texts = [
                cell.text.strip() for cell in row.cells
            ]
            texts.append(" | ".join(row_texts))

    # Group shapes (recursive)
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for child in shape.shapes:
            texts.append(_extract_text_from_shape(child))

    return "\n".join(texts)


# ---------------------------------------------------------------------------
# Adapter
# ---------------------------------------------------------------------------


@register_adapter
class EnterprisePPTAdapter(DataSourceAdapter):
    source_type = "enterprise_ppt"

    # ------------------------------------------------------------------
    # supports
    # ------------------------------------------------------------------

    def supports(self, filename: str) -> bool:
        return Path(filename).suffix.lower() == ".pptx"

    # ------------------------------------------------------------------
    # parse
    # ------------------------------------------------------------------

    async def parse(
        self,
        file_path: str,
        filename: str,
        metadata: Optional[dict] = None,
    ) -> SourceDocument:
        deck = await self.parse_deck(file_path, filename)
        md = self._deck_to_markdown(deck)

        return SourceDocument(
            id=deck.id,
            title=deck.title,
            content=md,
            source_type=self.source_type,
            metadata=metadata or {},
        )

    # ------------------------------------------------------------------
    # parse_deck
    # ------------------------------------------------------------------

    async def parse_deck(self, file_path: str, filename: str) -> ContentDeck:
        """Parse a .pptx file into a ContentDeck with slide summaries."""
        from pptx import Presentation  # type: ignore[import-untyped]

        prs = Presentation(file_path)
        deck_id = str(uuid.uuid4())
        title = Path(filename).stem

        slides: list[dict] = []
        slide_assets: list[SlideAsset] = []

        # Try to get LLM client (may fail if API key not set)
        try:
            from backend.core.llm.client import get_llm_client
            from backend.core.llm.prompts import SLIDE_SUMMARY_PROMPT
            llm = get_llm_client()
            llm_available = True
        except Exception:
            llm_available = False
            logger.warning("LLM client unavailable; using rule-based fallback for slide summaries")

        for i, slide in enumerate(prs.slides):
            # Extract all text
            slide_text_parts: list[str] = []
            for shape in slide.shapes:
                shape_text = _extract_text_from_shape(shape)
                if shape_text:
                    slide_text_parts.append(shape_text)

            full_text = "\n".join(slide_text_parts)

            # Determine slide title (first text paragraph, or fallback)
            slide_title = self._extract_slide_title(slide, i)

            # Summarize via LLM or fallback
            if llm_available and full_text.strip():
                try:
                    summary_json = await llm.chat_json(
                        system_prompt=SLIDE_SUMMARY_PROMPT.format(
                            slide_index=i + 1,
                            title=slide_title,
                            text_content=full_text[:3000],
                        ),
                        user_message="请分析以上幻灯片内容并输出 JSON。",
                        temperature=0.2,
                    )
                    summary = summary_json.get("summary", "")
                    keywords = summary_json.get("keywords", [])
                    slide_type = summary_json.get("slide_type", "content")
                except Exception:
                    summary = self._fallback_summary(full_text)
                    keywords = []
                    slide_type = "content"
            else:
                summary = self._fallback_summary(full_text)
                keywords = []
                slide_type = "content"

            # Build SlideAsset
            asset = SlideAsset(
                id=str(uuid.uuid4()),
                title=slide_title,
                description=summary,
                category=slide_type,
                tags=keywords,
                source_type="uploaded",
                file_path=file_path,
                slide_index=i,
                metadata={"slide_number": i + 1},
            )
            slide_assets.append(asset)

            # Build slide dict for ContentDeck
            slides.append({
                "index": i,
                "title": slide_title,
                "summary": summary,
                "keywords": keywords,
                "type": slide_type,
                "text": full_text,
                "slide_asset_id": asset.id,
            })

        # Extract style profile
        style = self.extract_style_profile(file_path, filename, prs)

        outline = {
            "total_slides": len(slides),
            "slide_titles": [s["title"] for s in slides],
        }

        return ContentDeck(
            id=deck_id,
            title=title,
            slides=slides,
            outline=outline,
        )

    # ------------------------------------------------------------------
    # extract_style_profile
    # ------------------------------------------------------------------

    def extract_style_profile(
        self,
        file_path: str,
        filename: str,
        prs=None,
    ) -> StyleProfile:
        """Extract basic font and layout info from slide masters."""
        from pptx import Presentation  # type: ignore[import-untyped]

        if prs is None:
            prs = Presentation(file_path)

        font_heading = "Arial"
        font_body = "Arial"
        color_palette: list[str] = []

        try:
            for master in prs.slide_masters:
                # Extract theme colors from slide master
                try:
                    theme = master.slide_layouts[0] if master.slide_layouts else None
                except Exception:
                    theme = None

                # Check layout placeholders for font info and fills
                for layout in master.slide_layouts:
                    for ph in layout.placeholders:
                        if ph.has_text_frame:
                            for para in ph.text_frame.paragraphs:
                                for run in para.runs:
                                    fname = _safe_get(run.font, "name")
                                    if fname:
                                        if "标题" in _safe_get(ph, "name", "") or "Title" in _safe_get(ph, "name", ""):
                                            font_heading = fname
                                        else:
                                            font_body = fname
                                        break
                            # Extract fill colors from shapes
                            try:
                                fill = ph.fill
                                if fill and fill.type is not None:
                                    try:
                                        fc = fill.fore_color
                                        if fc and fc.type is not None:
                                            rgb = str(fc.rgb) if hasattr(fc, 'rgb') else None
                                            if rgb and rgb not in color_palette:
                                                color_palette.append(f"#{rgb}")
                                    except Exception:
                                        pass
                            except Exception:
                                pass
                    if len(color_palette) >= 6:
                        break
                if len(color_palette) >= 3:
                    break
        except Exception:
            logger.warning("Could not extract style info from slide masters; using defaults")

        if not color_palette:
            color_palette = ["#003366", "#0066CC", "#FFFFFF", "#333333"]
        color_palette = color_palette[:6]

        # Extract slide dimensions
        w_emu = _safe_get(prs, "slide_width", 12192000)  # default 13.333" in EMU
        h_emu = _safe_get(prs, "slide_height", 6858000)   # default 7.5" in EMU

        return StyleProfile(
            id=str(uuid.uuid4()),
            name=Path(filename).stem,
            color_palette=color_palette or ["#1A1A1A", "#FFFFFF", "#007ACC"],
            font_heading=font_heading or "Arial",
            font_body=font_body or "Arial",
            slide_width_inches=round(w_emu / 914400, 3),   # EMU to inches
            slide_height_inches=round(h_emu / 914400, 3),
        )

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _extract_slide_title(slide, index: int) -> str:
        """Extract a title string from a slide."""
        # Try the slide title placeholder first
        try:
            title_shape = slide.shapes.title
            if title_shape and title_shape.has_text_frame:
                t = title_shape.text_frame.text.strip()
                if t:
                    return t
        except Exception:
            pass

        # Try first non-empty text shape
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    return t[:100]

        return f"Slide {index + 1}"

    @staticmethod
    def _fallback_summary(text: str, max_len: int = 200) -> str:
        """Rule-based fallback when LLM is unavailable."""
        if not text.strip():
            return "(empty slide)"
        lines = text.strip().split("\n")
        # Take first meaningful line(s) up to max_len
        result = lines[0].strip()
        if len(result) < 50 and len(lines) > 1:
            result += " " + lines[1].strip()
        return result[:max_len]

    @staticmethod
    def _deck_to_markdown(deck: ContentDeck) -> str:
        """Format a ContentDeck as structured Markdown."""
        parts: list[str] = []
        parts.append(f"# {deck.title}")
        parts.append(f"**Total Slides**: {deck.outline.get('total_slides', len(deck.slides))}")
        parts.append("")

        for slide in deck.slides:
            idx = slide["index"] + 1
            title = slide.get("title", f"Slide {idx}")
            summary = slide.get("summary", "")
            stype = slide.get("type", "content")
            keywords = slide.get("keywords", [])

            parts.append(f"## Slide {idx}: {title}")
            parts.append(f"**Type**: {stype}")
            if keywords:
                parts.append(f"**Keywords**: {', '.join(keywords)}")
            if summary:
                parts.append(f"**Summary**: {summary}")
            parts.append("")

        return "\n".join(parts)
