"""PPTX exporter — StructuredReport → native editable .pptx via SVG→PPTX pipeline."""
import os
import sys
import uuid
import logging
from pathlib import Path
from copy import deepcopy

from backend.core.output.base import BaseExporter
from backend.core.models import StructuredReport, ExportResult

logger = logging.getLogger(__name__)

# Ensure ppt-master scripts are importable
_PPT_MASTER_PATH = Path(__file__).resolve().parents[2] / "skills" / "ppt-master" / "scripts"
if str(_PPT_MASTER_PATH) not in sys.path:
    sys.path.insert(0, str(_PPT_MASTER_PATH))

# ppt169 canvas: 1280 x 720
CANVAS_W, CANVAS_H = 1280, 720

# Default color scheme
DEFAULT_COLORS = {
    "primary": "#003366",
    "accent": "#0066CC",
    "text": "#333333",
    "text_light": "#666666",
    "bg": "#FFFFFF",
    "bg_alt": "#F5F7FA",
}
DEFAULT_FONT = "Microsoft YaHei, SimHei, sans-serif"


class PptxExporter(BaseExporter):
    """Export StructuredReport to native editable PPTX via SVG→PPTX pipeline."""

    format = "pptx"

    async def export(
        self,
        report: StructuredReport,
        output_dir: str = "./data/exports",
        ppt_template_id: str = None,
    ) -> ExportResult:
        Path(output_dir).mkdir(parents=True, exist_ok=True)

        # 1. Generate SVG files from report sections
        svg_dir = Path(output_dir) / f"svg_{report.meta.id}"
        svg_dir.mkdir(parents=True, exist_ok=True)
        svg_files, notes = self._generate_svgs(report, svg_dir)

        if not svg_files:
            raise ValueError("No SVG files generated — report has no sections")

        # 2. SVG → PPTX
        output_path = Path(output_dir) / f"{report.meta.id}.pptx"
        from svg_to_pptx import create_pptx_with_native_svg

        success = create_pptx_with_native_svg(
            svg_files=svg_files,
            output_path=output_path,
            canvas_format="ppt169",
            verbose=False,
            use_native_shapes=True,
            notes=notes,
            enable_notes=True,
        )

        if not success:
            raise RuntimeError("PPTX creation failed — create_pptx_with_native_svg returned False")

        # 3. Handle slide reuse from enterprise PPT library (if any)
        self._merge_reused_slides(report, output_path)

        file_size = output_path.stat().st_size
        return ExportResult(
            file_path=str(output_path),
            format="pptx",
            file_size_bytes=file_size,
        )

    # ------------------------------------------------------------------
    # SVG generation
    # ------------------------------------------------------------------

    def _generate_svgs(self, report: StructuredReport, svg_dir: Path) -> tuple:
        """Generate SVG files for each report section. Returns (svg_paths, notes_dict)."""
        svg_files = []
        notes = {}
        sections = report.sections

        # Extract style from template if available
        colors = DEFAULT_COLORS
        font = DEFAULT_FONT

        for i, section in enumerate(sections):
            is_first = (i == 0)
            is_last = (i == len(sections) - 1)
            section_type = self._section_type(i, len(sections))

            svg_content = self._render_section_svg(
                section=section,
                report_title=report.meta.title,
                section_type=section_type,
                colors=colors,
                font=font,
                slide_num=i + 1,
                total=len(sections),
            )

            stem = f"slide_{i + 1:02d}"
            svg_path = svg_dir / f"{stem}.svg"
            svg_path.write_text(svg_content, encoding="utf-8")
            svg_files.append(svg_path)
            notes[stem] = section.markdown_content[:500] if section.markdown_content else section.section_def.title

        return svg_files, notes

    def _section_type(self, index: int, total: int) -> str:
        if index == 0:
            return "cover"
        if index == total - 1:
            return "ending"
        return "content"

    def _render_section_svg(self, section, report_title: str, section_type: str,
                            colors: dict, font: str, slide_num: int, total: int) -> str:
        """Render a single section as an SVG slide."""
        title = section.section_def.title
        content = section.markdown_content or ""

        # Parse content lines
        lines = [l.strip() for l in content.split("\n") if l.strip()]
        body_lines = []
        for line in lines:
            if line.startswith("#"):
                continue
            if line.startswith("- ") or line.startswith("* "):
                body_lines.append(("bullet", line[2:]))
            elif line.startswith("|"):
                body_lines.append(("table", line))
            else:
                body_lines.append(("text", line))

        # Limit lines to fit slide
        body_lines = body_lines[:20]

        if section_type == "cover":
            return self._cover_svg(title=report_title, subtitle=title, colors=colors, font=font)
        else:
            return self._content_svg(title, body_lines, colors=colors, font=font,
                                     slide_num=slide_num, total=total, is_ending=(section_type == "ending"))

    def _cover_svg(self, title: str = "", subtitle: str = "", description: str = "",
                   colors: dict = None, font: str = DEFAULT_FONT) -> str:
        colors = colors or DEFAULT_COLORS
        return f'''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {CANVAS_W} {CANVAS_H}">
  <defs>
    <linearGradient id="bg" x1="0" y1="0" x2="1" y2="1">
      <stop offset="0%" stop-color="{colors['primary']}"/>
      <stop offset="100%" stop-color="{colors['accent']}"/>
    </linearGradient>
  </defs>
  <rect width="{CANVAS_W}" height="{CANVAS_H}" fill="url(#bg)"/>
  <rect x="60" y="220" width="{CANVAS_W - 120}" height="280" rx="12" fill="rgba(255,255,255,0.1)"/>
  <text x="{CANVAS_W // 2}" y="320" text-anchor="middle" font-family="{font}" font-size="52" font-weight="bold" fill="white">{self._escape(title)}</text>
  <text x="{CANVAS_W // 2}" y="390" text-anchor="middle" font-family="{font}" font-size="24" fill="rgba(255,255,255,0.85)">{self._escape(subtitle)}</text>
  <text x="{CANVAS_W // 2}" y="440" text-anchor="middle" font-family="{font}" font-size="18" fill="rgba(255,255,255,0.7)">{self._escape(description)}</text>
  <line x1="260" y1="490" x2="{CANVAS_W - 260}" y2="490" stroke="rgba(255,255,255,0.4)" stroke-width="2"/>
  <text x="{CANVAS_W // 2}" y="540" text-anchor="middle" font-family="{font}" font-size="16" fill="rgba(255,255,255,0.6)">智能报告平台 · 自动生成</text>
</svg>'''

    def _content_svg(self, title: str, body_lines: list, colors: dict = None,
                     font: str = DEFAULT_FONT, slide_num: int = 1, total: int = 1,
                     is_ending: bool = False) -> str:
        colors = colors or DEFAULT_COLORS
        y = 120
        svg_parts = [
            f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {CANVAS_W} {CANVAS_H}">',
            f'<rect width="{CANVAS_W}" height="{CANVAS_H}" fill="{colors["bg"]}"/>',
            # Header bar
            f'<rect x="0" y="0" width="{CANVAS_W}" height="80" fill="{colors["primary"]}"/>',
            f'<text x="60" y="50" font-family="{font}" font-size="28" font-weight="bold" fill="white">{self._escape(title)}</text>',
            # Page number
            f'<text x="{CANVAS_W - 60}" y="50" text-anchor="end" font-family="{font}" font-size="14" fill="rgba(255,255,255,0.6)">{slide_num} / {total}</text>',
        ]

        for line_type, text in body_lines:
            if y > CANVAS_H - 60:
                break
            if line_type in ("text", "bullet"):
                prefix = "• " if line_type == "bullet" else ""
                svg_parts.append(
                    f'<text x="60" y="{y}" font-family="{font}" font-size="18" fill="{colors["text"]}">{self._escape(prefix + text[:120])}</text>'
                )
                y += 30
            elif line_type == "table":
                svg_parts.append(
                    f'<text x="60" y="{y}" font-family="{font}" font-size="14" fill="{colors["text_light"]}">{self._escape(text[:100])}</text>'
                )
                y += 24

        # Footer
        svg_parts.append(
            f'<line x1="60" y1="{CANVAS_H - 40}" x2="{CANVAS_W - 60}" y2="{CANVAS_H - 40}" stroke="{colors["bg_alt"]}" stroke-width="1"/>'
        )
        svg_parts.append(
            f'<text x="60" y="{CANVAS_H - 15}" font-family="{font}" font-size="11" fill="{colors["text_light"]}">智能报告平台</text>'
        )

        svg_parts.append('</svg>')
        return "\n".join(svg_parts)

    def _escape(self, text: str) -> str:
        """Escape text for SVG XML."""
        return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;")

    # ------------------------------------------------------------------
    # Slide reuse from enterprise PPT library
    # ------------------------------------------------------------------

    def _merge_reused_slides(self, report: StructuredReport, output_path: Path):
        """Replace placeholder slides with actual slides from enterprise PPT library.

        Strategy: For sections with accepted slide_refs pointing to enterprise PPT source
        files, find the source PPTX, clone the referenced slides, and insert them at the
        correct positions in the output PPTX, replacing the placeholder SVG slides.
        """
        from pptx import Presentation

        if not report.sections:
            return

        # Build a map: output_slide_index -> source_file -> source_slide_index
        # The output PPTX has slides in the same order as report.sections
        slide_replacements: list[tuple[int, str, int, str]] = []  # (output_idx, source_file, source_idx, slide_id)
        for i, section in enumerate(report.sections):
            for ref in section.slide_refs:
                if ref.accepted and ref.source_file and os.path.exists(ref.source_file):
                    slide_replacements.append((i, ref.source_file, ref.slide_index, ref.slide_id))

        if not slide_replacements:
            return

        try:
            prs = Presentation(str(output_path))
            total_slides = len(prs.slides)

            # Process replacements in reverse order to preserve indices
            for out_idx, source_file, source_idx, slide_id in sorted(slide_replacements, key=lambda x: x[0], reverse=True):
                if out_idx >= total_slides:
                    continue

                try:
                    src_prs = Presentation(source_file)
                    if source_idx >= len(src_prs.slides):
                        logger.warning(f"Source slide index {source_idx} out of range for {source_file}")
                        continue

                    # Clone slide from source into output at the correct position
                    self._clone_slide_into(src_prs, source_idx, prs, out_idx)

                except Exception as e:
                    logger.warning(f"Failed to clone slide {source_idx} from {source_file}: {e}")

            prs.save(str(output_path))
            logger.info(f"Slide reuse: {len(slide_replacements)} slides cloned from enterprise PPT library")

        except Exception as e:
            logger.warning(f"Slide reuse merge failed: {e}")

    def _clone_slide_into(self, src_prs, src_idx: int, dst_prs, dst_idx: int):
        """Clone a single slide from src_prs into dst_prs at position dst_idx."""
        import copy
        from lxml import etree

        src_slide = src_prs.slides[src_idx]

        # Create a new slide in destination with the same layout (or first available)
        try:
            layout_name = src_slide.slide_layout.name if src_slide.slide_layout else ""
            matching_layout = None
            for layout in dst_prs.slide_layouts:
                if layout.name == layout_name:
                    matching_layout = layout
                    break
            layout = matching_layout or dst_prs.slide_layouts[0]
        except Exception:
            layout = dst_prs.slide_layouts[0]

        # Add slide at the end first, then we'll reorder
        new_slide = dst_prs.slides.add_slide(layout)

        # Copy all shapes from source slide
        shape_tree = new_slide.shapes._spTree
        for shape in src_slide.shapes:
            try:
                el = copy.deepcopy(shape._element)
                shape_tree.insert_element_before(el, 'p:ext')
            except Exception:
                pass  # skip problematic shapes

        # Copy image relationships
        try:
            for rel in src_slide.part.rels.values():
                if "image" in rel.reltype or "media" in rel.reltype:
                    try:
                        new_slide.part.rels.get_or_add(rel.reltype, rel.target)
                    except Exception:
                        pass
        except Exception:
            pass

        # Move the newly added slide (at the end) to dst_idx
        total = len(dst_prs.slides)
        if total > 1 and total - 1 != dst_idx:
            try:
                xml_slides = dst_prs.slides._sldIdLst
                slide_entries = list(xml_slides)
                last_entry = slide_entries[-1]
                xml_slides.remove(last_entry)
                if dst_idx < len(slide_entries) - 1:
                    xml_slides.insert(dst_idx, last_entry)
                else:
                    xml_slides.append(last_entry)
            except Exception:
                pass  # best-effort reordering

    # ------------------------------------------------------------------
    # Speaker notes
    # ------------------------------------------------------------------

    def _build_notes(self, report: StructuredReport) -> dict[str, str]:
        """Build notes dict keyed by SVG filename stem."""
        notes = {}
        for i, section in enumerate(report.sections):
            stem = f"slide_{i + 1:02d}"
            content = section.markdown_content or ""
            notes[stem] = content[:1000]
        return notes
