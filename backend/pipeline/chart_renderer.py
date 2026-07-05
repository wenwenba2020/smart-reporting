"""W2-4: 图表渲染器 — 输出 python-pptx ChartData 用于 add_chart() 插入原生可编辑图表。

规范要求（agents-spec.md / known-pitfalls.md 坑10）：
  图表必须用 python-pptx add_chart() 插入 DrawingML，不是插入 SVG 图片。
  SVG 插入后会变成光栅图，不可编辑数据。

本模块输出 ChartDefinition（chart_type + ChartData），由编辑师调用 slide.shapes.add_chart() 插入。
"""
from dataclasses import dataclass, field
from typing import Any

from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE

# Default color palette (hex, applied via python-pptx after chart creation)
DEFAULT_COLORS = ["#4A90D9", "#F5A623", "#5DCAA5", "#F0997B", "#8B5CF6"]

# Map from our chart type names to XL_CHART_TYPE enums
CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "bar-grouped": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "radar": XL_CHART_TYPE.RADAR,
    "funnel": XL_CHART_TYPE.BAR_STACKED,  # no native funnel in Office, use stacked bar
}


@dataclass
class ChartDefinition:
    """Ready-to-insert chart definition for python-pptx add_chart()."""
    xl_chart_type: int
    chart_data: ChartData
    colors: list[str] = field(default_factory=lambda: list(DEFAULT_COLORS))
    caption: str = ""


def build_category_chart(
    chart_type: str,
    categories: list[str],
    series: list[dict[str, Any]],
    caption: str = "",
    colors: list[str] | None = None,
) -> ChartDefinition:
    """Build a category-based chart (bar, bar-grouped, line, radar)."""
    xl_type = CHART_TYPE_MAP.get(chart_type)
    if xl_type is None:
        raise ValueError(f"Unsupported chart type: {chart_type}")

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series:
        chart_data.add_series(s["name"], s["values"])

    return ChartDefinition(
        xl_chart_type=xl_type,
        chart_data=chart_data,
        colors=colors or list(DEFAULT_COLORS),
        caption=caption,
    )


def build_pie_chart(
    data: list[tuple[str, float] | list],
    caption: str = "",
    colors: list[str] | None = None,
) -> ChartDefinition:
    """Build a pie chart from (label, value) pairs."""
    chart_data = CategoryChartData()
    labels = [d[0] for d in data]
    values = [d[1] for d in data]
    chart_data.categories = labels
    chart_data.add_series("", values)

    return ChartDefinition(
        xl_chart_type=XL_CHART_TYPE.PIE,
        chart_data=chart_data,
        colors=colors or list(DEFAULT_COLORS),
        caption=caption,
    )


def build_scatter_chart(
    data: list[list[float]],
    caption: str = "",
    colors: list[str] | None = None,
) -> ChartDefinition:
    """Build a scatter chart from [[x, y], ...] data."""
    from pptx.chart.data import XyChartData

    chart_data = XyChartData()
    series = chart_data.add_series("")
    for point in data:
        series.add_data_point(point[0], point[1])

    return ChartDefinition(
        xl_chart_type=XL_CHART_TYPE.XY_SCATTER,
        chart_data=chart_data,
        colors=colors or list(DEFAULT_COLORS),
        caption=caption,
    )


# Dispatcher: chart_type → builder function
BUILDERS = {
    "bar": lambda cfg, colors: build_category_chart(
        "bar", cfg["categories"], cfg["series"], cfg.get("caption", ""), colors,
    ),
    "bar-grouped": lambda cfg, colors: build_category_chart(
        "bar-grouped", cfg["categories"], cfg["series"], cfg.get("caption", ""), colors,
    ),
    "line": lambda cfg, colors: build_category_chart(
        "line", cfg["categories"], cfg["series"], cfg.get("caption", ""), colors,
    ),
    "pie": lambda cfg, colors: build_pie_chart(
        cfg["data"], cfg.get("caption", ""), colors,
    ),
    "scatter": lambda cfg, colors: build_scatter_chart(
        cfg["data"], cfg.get("caption", ""), colors,
    ),
    "radar": lambda cfg, colors: build_category_chart(
        "radar", cfg["categories"], cfg["series"], cfg.get("caption", ""), colors,
    ),
    "funnel": lambda cfg, colors: build_pie_chart(
        cfg["data"], cfg.get("caption", ""), colors,
    ),
}


def render_chart(
    chart_type: str,
    chart_config: dict[str, Any],
    colors: list[str] | None = None,
) -> ChartDefinition:
    """
    Build a ChartDefinition ready for python-pptx add_chart().

    Usage:
        defn = render_chart("bar", {"categories": [...], "series": [...]})
        slide.shapes.add_chart(
            defn.xl_chart_type, left, top, width, height, defn.chart_data
        )
    """
    builder = BUILDERS.get(chart_type)
    if builder is None:
        raise ValueError(f"Unsupported chart type: {chart_type}. Supported: {list(BUILDERS.keys())}")
    return builder(chart_config, colors)
