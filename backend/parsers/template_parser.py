"""企业报告模板解析器 — 从 PPTX/DOCX 提取样式规则和内容插槽"""
import re
from pathlib import Path


def parse_pptx_template(file_path: str | Path) -> dict:
    """从 PPTX 模板提取样式规则 + 内容插槽。"""
    from pptx import Presentation

    file_path = Path(file_path)
    prs = Presentation(str(file_path))

    colors: set[str] = set()
    fonts: set[str] = set()
    font_sizes: set[int] = set()
    slots: list[dict] = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        text = run.text.strip()
                        if run.font.color and run.font.color.rgb:
                            colors.add(str(run.font.color.rgb))
                        if run.font.name:
                            fonts.add(run.font.name)
                        if run.font.size:
                            font_sizes.add(int(run.font.size.pt))

                        placeholders = re.findall(r'\{\{(\w+)\}\}', text)
                        for ph in placeholders:
                            slots.append({
                                "slot_id": ph,
                                "type": _infer_slot_type(ph),
                                "slide_index": slide_idx,
                                "shape_name": shape.name,
                            })

    return {
        "style_rules": {
            "colors": sorted(colors)[:12],
            "fonts": sorted(fonts),
            "font_sizes": sorted(font_sizes),
        },
        "content_slots": slots,
        "slide_count": len(prs.slides),
    }


def parse_docx_template(file_path: str | Path) -> dict:
    """从 DOCX 模板提取样式规则 + 内容插槽。"""
    from docx import Document

    file_path = Path(file_path)
    doc = Document(str(file_path))

    fonts: set[str] = set()
    font_sizes: set[int] = set()
    slots: list[dict] = []

    for para in doc.paragraphs:
        for run in para.runs:
            if run.font.name:
                fonts.add(run.font.name)
            if run.font.size:
                font_sizes.add(int(run.font.size.pt))

            placeholders = re.findall(r'\{\{(\w+)\}\}', run.text)
            for ph in placeholders:
                slots.append({
                    "slot_id": ph,
                    "type": _infer_slot_type(ph),
                    "paragraph_style": para.style.name if para.style else None,
                })

    return {
        "style_rules": {"fonts": sorted(fonts), "font_sizes": sorted(font_sizes)},
        "content_slots": slots,
    }


def _infer_slot_type(slot_id: str) -> str:
    """根据占位符名称推断内容类型。"""
    slot_lower = slot_id.lower()
    if any(k in slot_lower for k in ("title", "标题", "heading")):
        return "title"
    if any(k in slot_lower for k in ("chart", "图表", "graph")):
        return "chart"
    if any(k in slot_lower for k in ("table", "表格", "grid")):
        return "table"
    if any(k in slot_lower for k in ("image", "图片", "logo", "photo")):
        return "image"
    if any(k in slot_lower for k in ("date", "日期", "author", "作者")):
        return "metadata"
    return "text"
