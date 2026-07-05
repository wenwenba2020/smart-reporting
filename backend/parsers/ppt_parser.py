"""W2-2: 参考 PPT 解析 — 图片提取 + 配色/字体信息"""
import hashlib
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


def extract_images(pptx_path: str | Path, output_dir: str | Path) -> list[str]:
    """Extract all images from PPTX, save to output_dir. Returns list of saved paths."""
    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))
    saved: list[str] = []
    seen_blobs: set[str] = set()  # dedupe by SHA-256

    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.shape_type or shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            image = shape.image
            blob = image.blob
            blob_hash = hashlib.sha256(blob).hexdigest()
            if blob_hash in seen_blobs:
                continue
            seen_blobs.add(blob_hash)

            ext = image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            filename = f"slide_{slide_idx:02d}_img_{len(saved) + 1}.{ext}"
            out_path = output_dir / filename
            out_path.write_bytes(blob)
            saved.append(str(out_path))

    return saved


def extract_style_info(pptx_path: str | Path) -> dict:
    """Extract color scheme and font info from PPTX. Returns JSON-serializable dict."""
    prs = Presentation(str(pptx_path))

    colors: set[str] = set()
    fonts: set[str] = set()
    font_sizes: set[float] = set()

    for slide in prs.slides:
        # Slide background color
        bg = slide.background
        if bg.fill and bg.fill.type is not None:
            try:
                fc = bg.fill.fore_color
                if fc and fc.type is not None:
                    colors.add(str(fc.rgb))
            except (AttributeError, TypeError):
                pass

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        fonts.add(run.font.name)
                    if run.font.size:
                        font_sizes.add(run.font.size.pt)
                    if run.font.color and run.font.color.type is not None:
                        try:
                            colors.add(str(run.font.color.rgb))
                        except (AttributeError, TypeError):
                            pass

    return {
        "colors": sorted(colors),
        "fonts": sorted(fonts),
        "font_sizes": sorted(font_sizes),
        "slide_count": len(prs.slides),
    }


def analyze_reference_pptx(
    pptx_path: str | Path,
    output_dir: str | Path,
) -> dict:
    """Full analysis: extract images + style info, save JSON description."""
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    images = extract_images(pptx_path, output_dir)
    style = extract_style_info(pptx_path)

    result = {
        "source": str(pptx_path),
        "images": images,
        "style": style,
    }

    desc_path = output_dir / "reference_analysis.json"
    desc_path.write_text(json.dumps(result, ensure_ascii=False, indent=2))
    return result
