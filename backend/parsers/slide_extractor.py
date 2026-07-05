"""从 PPTX 文件逐页提取 slide 内容、XML、图片、缩略图"""
import re
import zipfile
from pathlib import Path
from xml.sax.saxutils import escape as xml_escape
from dataclasses import dataclass, field

from pptx import Presentation


@dataclass
class ExtractedSlide:
    slide_index: int          # 1-based
    title: str                # first text frame content (truncated)
    text_summary: str         # all text content concatenated
    slide_xml_rel: str        # relative path: {lib_id}/slides/slide_NN/slide.xml
    thumbnail_rel: str        # relative path: {lib_id}/slides/slide_NN/thumbnail.svg
    layout_hint: str          # guessed layout type
    tags: list[str] = field(default_factory=list)


SIMPLE_LAYOUT_RULES = [
    (re.compile(r"(目录|contents|toc|agenda|outline)", re.I), "toc"),
    (re.compile(r"(团队|team|成员|member|组织)", re.I), "team"),
    (re.compile(r"(对比|比较|vs|comparison|versus)", re.I), "comparison"),
    (re.compile(r"(时间|timeline|历程|历史|里程碑|milestone)", re.I), "timeline"),
    (re.compile(r"(图表|数据|chart|graph|statistics|指标)", re.I), "data-chart"),
    (re.compile(r"(引用|quote|名言|证言|testimonial)", re.I), "quote"),
]


def _guess_layout(title_text: str, body_text: str, slide_index: int, total: int) -> str:
    """Guess layout type from text content."""
    combined = f"{title_text} {body_text[:200]}"
    if slide_index == 1 and total > 3:
        return "cover"
    for pattern, layout in SIMPLE_LAYOUT_RULES:
        if pattern.search(combined):
            return layout
    return "title-content"


def generate_ai_tags(title: str, text_summary: str, existing_hint: str = "") -> list[str]:
    """Generate 2-4 Chinese tags for a slide using a lightweight LLM. Falls back gracefully on error."""
    prompt = f"""根据以下幻灯片内容，生成 2-4 个中文标签（逗号分隔）。
标签应简短（2-6 字），描述页面类型和主题。
常见标签：封面、目录、团队介绍、数据图表、时间线、对比分析、引用名言、总结、产品介绍、流程图、SWOT分析、用户画像、商业模式、路线图、KPI指标

标题：{title or '(无)'}
内容：{(text_summary or '')[:500]}

只返回逗号分隔的标签，不要其他文字。"""

    try:
        from backend.agents.llm_client import chat
        result = chat(
            model="qwen/qwen3.5-9b",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3,
            max_tokens=64,
            extra_body={"enable_thinking": False},
        )
        tags = [t.strip() for t in re.split(r"[,，、]", result) if t.strip()]
        if existing_hint and existing_hint not in tags:
            tags.append(existing_hint)
        return tags[:6]
    except Exception:
        return [existing_hint] if existing_hint else []


EMU_PER_PX = 914400 / 96


def _render_svg_thumbnail(slide, width: int = 960, height: int = 540) -> str:
    """Generate SVG thumbnail for a single slide (text + approximate positions)."""
    bg_color = "#ffffff"
    try:
        bg = slide.background
        if bg.fill and bg.fill.type is not None:
            fc = bg.fill.fore_color
            if fc and fc.type is not None:
                bg_color = f"#{fc.rgb}"
    except Exception:
        pass

    lines: list[str] = [
        f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {width} {height}">',
        f'<rect width="{width}" height="{height}" fill="{bg_color}"/>',
    ]

    for shape in slide.shapes:
        try:
            left = shape.left / EMU_PER_PX if shape.left else 0
            top = shape.top / EMU_PER_PX if shape.top else 0
            w = shape.width / EMU_PER_PX if shape.width else width
            h = shape.height / EMU_PER_PX if shape.height else 20
            text = shape.text_frame.text.strip() if shape.has_text_frame else ""
        except Exception:
            continue

        if not text:
            # Image placeholder
            if shape.shape_type and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                lines.append(
                    f'<rect x="{left:.0f}" y="{top:.0f}" width="{w:.0f}" height="{h:.0f}" '
                    f'fill="#e5e7eb" stroke="#d1d5db" rx="4"/>'
                    f'<text x="{left + w/2:.0f}" y="{top + h/2:.0f}" '
                    f'text-anchor="middle" dominant-baseline="central" font-size="11" fill="#9ca3af">[img]</text>'
                )
            continue

        font_size = 12
        is_bold = False
        try:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size and run.font.size.pt:
                            font_size = run.font.size.pt
                        if run.font.bold:
                            is_bold = True
                        break
                    break
        except Exception:
            pass

        weight = "bold" if is_bold else "normal"
        # Split into lines for multi-line text visibility
        text_lines = text.split("\n")[:8]
        line_height = font_size + 4
        for li, line_text in enumerate(text_lines):
            y = top + line_height * (li + 1)
            display_text = line_text[:80] + ("..." if len(line_text) > 80 else "")
            lines.append(
                f'<text x="{left:.0f}" y="{y:.0f}" font-size="{font_size}" '
                f'font-weight="{weight}" fill="#1f2937" font-family="sans-serif">'
                f'{xml_escape(display_text)}</text>'
            )

    lines.append("</svg>")
    return "\n".join(lines)


def extract_slides_from_pptx(pptx_path: str | Path, output_dir: str | Path) -> list[ExtractedSlide]:
    """Extract slides page-by-page from a PPTX file. Returns list of ExtractedSlide."""
    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))
    total = len(prs.slides)
    results: list[ExtractedSlide] = []

    # Extract per-slide XML from the ZIP
    slides_xml: dict[int, bytes] = {}
    with zipfile.ZipFile(pptx_path, "r") as zf:
        for name in zf.namelist():
            m = re.match(r"ppt/slides/slide(\d+)\.xml", name)
            if m:
                slide_num = int(m.group(1))
                slides_xml[slide_num] = zf.read(name)

    for idx, slide in enumerate(prs.slides, 1):
        slide_dir = output_dir / f"slide_{idx:02d}"
        slide_dir.mkdir(parents=True, exist_ok=True)

        # Extract text
        all_texts: list[str] = []
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    all_texts.append(text)
                    if shape.shape_type == 1 and not title:  # PLACEHOLDER title shape
                        title = text
        if not title and all_texts:
            title = all_texts[0][:60]
        text_summary = "\n\n".join(all_texts)

        # Save raw slide XML
        slide_xml_path = slide_dir / "slide.xml"
        if idx in slides_xml:
            slide_xml_path.write_bytes(slides_xml[idx])

        # Generate SVG thumbnail
        thumbnail_path = slide_dir / "thumbnail.svg"
        svg = _render_svg_thumbnail(slide)
        thumbnail_path.write_text(svg, encoding="utf-8")

        # Extract images
        images_dir = slide_dir / "images"
        images_dir.mkdir(exist_ok=True)
        img_count = 0
        for shape in slide.shapes:
            if shape.shape_type and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image = shape.image
                    blob = image.blob
                    ext = image.content_type.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"
                    img_path = images_dir / f"img_{img_count + 1}.{ext}"
                    img_path.write_bytes(blob)
                    img_count += 1
                except Exception:
                    pass

        layout_hint = _guess_layout(title, text_summary, idx, total)

        # Generate AI tags
        tags = generate_ai_tags(title, text_summary, layout_hint)

        results.append(ExtractedSlide(
            slide_index=idx,
            title=title or f"Slide {idx}",
            text_summary=text_summary,
            slide_xml_rel=str((slide_dir / "slide.xml").relative_to(output_dir.parent)),
            thumbnail_rel=str(thumbnail_path.relative_to(output_dir.parent)),
            layout_hint=layout_hint,
            tags=tags,
        ))

    return results
