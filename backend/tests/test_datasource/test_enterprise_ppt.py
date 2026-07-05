import pytest
import tempfile
import os

from backend.core.datasource.enterprise_ppt import EnterprisePPTAdapter


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _create_minimal_pptx() -> str:
    """Create a minimal .pptx file and return its path."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1 — title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # title layout
    title = slide1.shapes.title
    title.text = "Q3 营收分析报告"
    if slide1.placeholders[1].has_text_frame:
        slide1.placeholders[1].text = "2024年第三季度"

    # Slide 2 — content slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # title+content layout
    title2 = slide2.shapes.title
    title2.text = "核心数据概览"
    body = slide2.placeholders[1]
    body.text = "总营收: 5200万元\n净利润: 800万元\n同比增长: 15%"

    prs.save(path)
    return path


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_enterprise_ppt_parse():
    adapter = EnterprisePPTAdapter()
    path = _create_minimal_pptx()
    try:
        doc = await adapter.parse(path, "Q3_report.pptx")
        assert doc.source_type == "enterprise_ppt"
        assert "Q3 营收分析报告" in doc.content
        assert "核心数据概览" in doc.content
        assert "Total Slides" in doc.content
        assert "Slide 1" in doc.content
        assert "Slide 2" in doc.content
    finally:
        os.unlink(path)


@pytest.mark.asyncio
async def test_enterprise_ppt_parse_deck():
    adapter = EnterprisePPTAdapter()
    path = _create_minimal_pptx()
    try:
        deck = await adapter.parse_deck(path, "Q3_report.pptx")
        assert deck.title == "Q3_report"
        assert len(deck.slides) == 2
        assert deck.slides[0]["title"] == "Q3 营收分析报告"
        assert deck.slides[1]["title"] == "核心数据概览"
        assert "total_slides" in deck.outline
        assert deck.outline["total_slides"] == 2
        # Each slide should have a summary
        assert deck.slides[0]["summary"]
        assert deck.slides[1]["summary"]
    finally:
        os.unlink(path)


@pytest.mark.asyncio
async def test_enterprise_ppt_style_profile():
    adapter = EnterprisePPTAdapter()
    path = _create_minimal_pptx()
    try:
        from pptx import Presentation
        prs = Presentation(path)
        style = adapter.extract_style_profile(path, "Q3_report.pptx", prs)
        assert style.name == "Q3_report"
        assert style.font_heading
        assert style.font_body
        assert style.slide_width_inches > 0
        assert style.slide_height_inches > 0
        assert len(style.color_palette) > 0
    finally:
        os.unlink(path)


@pytest.mark.asyncio
async def test_enterprise_ppt_empty_slide():
    """Empty slide should not crash."""
    from pptx import Presentation

    adapter = EnterprisePPTAdapter()

    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)

    try:
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        prs.save(path)

        deck = await adapter.parse_deck(path, "empty.pptx")
        assert len(deck.slides) == 1
        # Should have a fallback summary, not crash
        assert deck.slides[0]["title"] == "Slide 1"
        assert deck.slides[0]["summary"] == "(empty slide)"
    finally:
        os.unlink(path)


def test_enterprise_ppt_supports():
    adapter = EnterprisePPTAdapter()
    assert adapter.supports("report.pptx") is True
    assert adapter.supports("slides.PPTX") is True
    assert adapter.supports("document.docx") is False
    assert adapter.supports("sheet.xlsx") is False
    assert adapter.supports("notes.txt") is False
